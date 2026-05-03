"""
Industry-grade PPT generation service.
Features: entrance animations, styled charts with data labels,
professional layouts, varied transitions, numbered bullet badges.
"""

import os
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from lxml import etree

from app.services.template_service import get_template, rgb

# ── Canvas dimensions (16:9 widescreen) ──────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

CHART_TYPE_MAP = {
    'bar':    XL_CHART_TYPE.BAR_CLUSTERED,
    'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
    'line':   XL_CHART_TYPE.LINE_MARKERS,
    'pie':    XL_CHART_TYPE.PIE,
    'area':   XL_CHART_TYPE.AREA,
}

# Palette used for multi-series charts and alternating bullet numbers
_CHART_PALETTE = [
    '#4F46E5', '#06B6D4', '#10B981', '#F59E0B',
    '#EF4444', '#8B5CF6', '#EC4899', '#14B8A6',
]


# ═══════════════════════ LOW-LEVEL HELPERS ═══════════════════════

def _solid_fill(shape, hex_color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb(hex_color)


def _no_line(shape):
    shape.line.fill.background()


def _set_font(run, size, bold=False, color=None, font_name='Calibri', italic=False):
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    if color:
        run.font.color.rgb = rgb(color)


def _add_rect(slide, left, top, width, height, hex_color, rounded=False):
    stype = 5 if rounded else 1
    shape = slide.shapes.add_shape(
        stype, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    _solid_fill(shape, hex_color)
    _no_line(shape)
    return shape


def _add_oval(slide, left, top, width, height, hex_color):
    shape = slide.shapes.add_shape(
        9,  # oval
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    _solid_fill(shape, hex_color)
    _no_line(shape)
    return shape


def _add_textbox(slide, text, left, top, width, height,
                 pt, bold, color, font_name='Calibri',
                 align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, pt, bold, color, font_name, italic)
    return tb


def _add_slide_number(slide, num, total, tmpl):
    """Bottom-right footer: slide X / Y."""
    tb = slide.shapes.add_textbox(
        Inches(11.5), Inches(7.1), Inches(1.6), Inches(0.3)
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f'{num}  /  {total}'
    _set_font(run, Pt(10), False, tmpl['subtitle_color'], tmpl['font_name'])


def _add_number_badge(slide, number, left, top, primary_hex, text_hex, size=0.38):
    """Circular number badge for bullet points."""
    badge = _add_oval(slide, left, top, size, size, primary_hex)
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(size), Inches(size)
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    _set_font(run, Pt(12), True, text_hex)
    return badge, tb


# ═══════════════════════ TRANSITIONS ═════════════════════════════

_TRANSITIONS = {
    'fade':  '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med" advTm="0"><p:fade/></p:transition>',
    'push':  '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med" advTm="0"><p:push dir="l"/></p:transition>',
    'wipe':  '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med" advTm="0"><p:wipe dir="l"/></p:transition>',
    'split': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med" advTm="0"><p:split orient="horz" dir="in"/></p:transition>',
    'zoom':  '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med" advTm="0"><p:zoom dir="in"/></p:transition>',
}


def _add_transition(slide_elem, trans_type='fade'):
    xml = _TRANSITIONS.get(trans_type, _TRANSITIONS['fade'])
    # Remove existing transition if any
    for child in list(slide_elem):
        if child.tag.endswith('}transition'):
            slide_elem.remove(child)
    slide_elem.append(etree.fromstring(xml))


# ═══════════════════════ ANIMATIONS ══════════════════════════════

def _build_animations(slide, anim_list):
    """
    Add entrance animations to a slide.
    anim_list: [(shape_id, delay_ms), ...]
    All use auto-play fade-in with staggered delays.
    """
    if not anim_list:
        return

    PML = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    cid = 3
    blocks = []
    bld_entries = []

    for grp, (spid, delay_ms) in enumerate(anim_list):
        i1 = cid;     cid += 1
        i2 = cid;     cid += 1
        i3 = cid;     cid += 1
        i4 = cid;     cid += 1

        blocks.append(
            f'<p:par>'
            f'<p:cTn id="{i1}" fill="hold">'
            f'<p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst>'
            f'<p:childTnLst>'
            f'<p:par><p:cTn id="{i2}" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>'
            f'<p:par><p:cTn id="{i3}" presetID="10" presetClass="entr" presetSubtype="0" '
            f'fill="hold" grpId="{grp}" nodeType="afterEffect">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>'
            f'<p:animEffect transition="in" filter="fade">'
            f'<p:cBhvr><p:cTn id="{i4}" dur="450"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
            f'</p:animEffect>'
            f'</p:childTnLst></p:cTn></p:par>'
            f'</p:childTnLst></p:cTn></p:par>'
            f'</p:childTnLst></p:cTn></p:par>'
        )
        bld_entries.append(
            f'<p:bldP spid="{spid}" grpId="{grp}" uiExpand="1" build="p"/>'
        )

    excl_id = cid
    timing_xml = (
        f'<p:timing xmlns:p="{PML}">'
        f'<p:tnLst><p:par>'
        f'<p:cTn id="1" dur="indefin" restart="whenNotActive" nodeType="tmRoot">'
        f'<p:childTnLst>'
        f'<p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="2" dur="indefin" nodeType="mainSeq">'
        f'<p:childTnLst>{"".join(blocks)}</p:childTnLst>'
        f'</p:cTn></p:seq>'
        f'<p:excl><p:cTn id="{excl_id}" dur="1000" fill="hold">'
        f'<p:stCondLst><p:cond evt="onNext" delay="0"><p:tn/></p:cond></p:stCondLst>'
        f'</p:cTn></p:excl>'
        f'</p:childTnLst></p:cTn></p:par></p:tnLst>'
        f'<p:bldLst>{"".join(bld_entries)}</p:bldLst>'
        f'</p:timing>'
    )

    try:
        # Remove existing timing if present
        for child in list(slide._element):
            if child.tag.endswith('}timing'):
                slide._element.remove(child)
        slide._element.append(etree.fromstring(timing_xml))
    except Exception:
        pass  # Never crash the PPT over animations


# ═══════════════════════ CHART STYLING ═══════════════════════════

def _set_series_color(series, hex_color):
    """Set series fill color via direct XML."""
    ser = series._element
    existing = ser.find(qn('c:spPr'))
    if existing is not None:
        ser.remove(existing)
    val = hex_color.lstrip('#')
    xml = (
        f'<c:spPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" '
        f'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:solidFill><a:srgbClr val="{val}"/></a:solidFill>'
        f'<a:ln><a:noFill/></a:ln>'
        f'<a:effectLst/>'
        f'</c:spPr>'
    )
    try:
        ser.append(etree.fromstring(xml))
    except Exception:
        pass


def _set_point_colors(series, palette):
    """
    Color each data point individually (per-bar coloring for single-series charts).
    Inserts <c:dPt> elements into the series XML with idx and spPr.
    """
    ser = series._element
    C = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
    A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # Find insertion point — dPt must come after order/tx/spPr
    # We'll collect the current child tags to find the right position
    insert_after = None
    for child in ser:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('order', 'tx', 'spPr', 'invertIfNegative', 'pictureOptions'):
            insert_after = child

    # Count existing data points by looking at val or cat
    count = 0
    val_elem = ser.find(f'{{{C}}}val')
    if val_elem is None:
        val_elem = ser.find(f'{{{C}}}yVal')
    if val_elem is not None:
        num_ref = val_elem.find(f'.//{{{C}}}numRef')
        if num_ref is not None:
            num_cache = num_ref.find(f'{{{C}}}numCache')
            if num_cache is not None:
                count = len(num_cache.findall(f'{{{C}}}pt'))

    if count == 0:
        return  # Can't determine point count; skip

    dpt_xmls = []
    for idx, color in zip(range(count), [palette[i % len(palette)] for i in range(count)]):
        val = color.lstrip('#')
        dpt_xml = (
            f'<c:dPt xmlns:c="{C}" xmlns:a="{A}">'
            f'<c:idx val="{idx}"/>'
            f'<c:invertIfNegative val="0"/>'
            f'<c:spPr>'
            f'<a:solidFill><a:srgbClr val="{val}"/></a:solidFill>'
            f'<a:ln><a:noFill/></a:ln>'
            f'<a:effectLst/>'
            f'</c:spPr>'
            f'</c:dPt>'
        )
        dpt_xmls.append(etree.fromstring(dpt_xml))

    # Insert dPt elements after insert_after (or at beginning if not found)
    if insert_after is not None:
        pos = list(ser).index(insert_after) + 1
        for i, dpt in enumerate(dpt_xmls):
            ser.insert(pos + i, dpt)
    else:
        for i, dpt in enumerate(dpt_xmls):
            ser.insert(i, dpt)


def _set_chart_bg(chart_obj, bg_hex):
    """Set chart plot area and chart area background color."""
    try:
        val = bg_hex.lstrip('#')
        bg_xml = (
            f'<c:spPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" '
            f'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:solidFill><a:srgbClr val="{val}"/></a:solidFill>'
            f'<a:ln><a:noFill/></a:ln>'
            f'</c:spPr>'
        )
        for elem_name in ('plotArea', 'chartSpace'):
            elem = chart_obj._element.find(
                f'.//{{{chart_obj._element.nsmap.get("c","http://schemas.openxmlformats.org/drawingml/2006/chart")}}}{elem_name}'
            )
            if elem is not None:
                existing = elem.find('{http://schemas.openxmlformats.org/drawingml/2006/chart}spPr')
                if existing is not None:
                    elem.remove(existing)
                elem.append(etree.fromstring(bg_xml))
    except Exception:
        pass


def _style_chart(chart_obj, tmpl, chart_type='bar'):
    """Professional chart styling: vibrant colors, data labels, clean axes."""
    try:
        chart_obj.has_legend = len(list(chart_obj.series)) > 1
        chart_obj.has_title = False
    except Exception:
        pass

    # Always start with the vivid chart palette so bars are readable on any bg
    palette = _CHART_PALETTE + [tmpl['accent'], tmpl['primary'], tmpl['secondary']]

    # Background: use template's own background (slightly off-white / dark)
    _set_chart_bg(chart_obj, tmpl['background'])

    # Series colors
    num_series = len(list(chart_obj.series))
    for i, series in enumerate(chart_obj.series):
        try:
            _set_series_color(series, palette[i % len(palette)])
            # For single-series bar/column, color each point individually
            if num_series == 1 and chart_type in ('bar', 'column'):
                _set_point_colors(series, palette)
        except Exception:
            pass

    # Data labels
    try:
        plot = chart_obj.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        try:
            dl.font.size = Pt(11)
            dl.font.bold = True
            # Use white for dark palette bars, dark for light ones
            dl.font.color.rgb = rgb('#FFFFFF')
        except Exception:
            pass
    except Exception:
        pass

    # Axis font styling
    for axis_name in ('category_axis', 'value_axis'):
        try:
            axis = getattr(chart_obj, axis_name)
            axis.tick_labels.font.size = Pt(10)
            axis.tick_labels.font.bold = False
            axis.tick_labels.font.color.rgb = rgb(tmpl['text_color'])
            # Thin axis lines
            axis.format.line.color.rgb = rgb(tmpl['subtitle_color'])
        except Exception:
            pass

    # Show legend only for multi-series; style it
    try:
        if chart_obj.has_legend:
            legend = chart_obj.legend
            legend.font.size = Pt(10)
            legend.font.color.rgb = rgb(tmpl['text_color'])
    except Exception:
        pass


def _add_chart(slide, chart_data_dict, left, top, width, height, tmpl):
    labels   = chart_data_dict.get('labels', ['A', 'B', 'C', 'D'])
    values   = chart_data_dict.get('values', [10, 20, 30, 40])
    ctype    = chart_data_dict.get('chart_type', 'bar')
    sname    = chart_data_dict.get('series_name', 'Value')

    xl_type  = CHART_TYPE_MAP.get(ctype, XL_CHART_TYPE.BAR_CLUSTERED)

    cd = ChartData()
    cd.categories = labels
    cd.add_series(sname, values)

    graphic_frame = slide.shapes.add_chart(
        xl_type,
        Inches(left), Inches(top), Inches(width), Inches(height),
        cd
    )
    chart_obj = graphic_frame.chart
    _style_chart(chart_obj, tmpl, ctype)
    return graphic_frame


# ═══════════════════════ SLIDE BUILDERS ══════════════════════════

def _make_title_slide(prs, slide_data, tmpl, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sp    = slide.shapes

    title_text    = slide_data.get('title', 'Presentation Title')
    subtitle_text = slide_data.get('subtitle', '')

    # ── Full background ──────────────────────────────────────────
    bg = _add_rect(slide, 0, 0, 13.33, 7.5, tmpl['primary'])

    # ── Decorative large circle (top-right, lighter shade) ───────
    # Use accent color at reduced visual weight via lighter hex
    circ1 = _add_oval(slide, 9.8, -1.8, 5.5, 5.5, tmpl['secondary'])
    circ2 = _add_oval(slide, 10.9, -0.6, 3.5, 3.5, tmpl['accent'])

    # ── Diagonal accent band (bottom-left) ───────────────────────
    band = _add_rect(slide, 0, 6.4, 13.33, 1.1, tmpl['secondary'])

    # ── Thin accent line across center ───────────────────────────
    line = _add_rect(slide, 0.55, 3.9, 4.0, 0.06, tmpl['accent'])

    # ── Title ─────────────────────────────────────────────────────
    tb_title = _add_textbox(
        slide, title_text,
        0.55, 1.4, 10.5, 2.6,
        Pt(48), True, '#FFFFFF', tmpl['font_name'], PP_ALIGN.LEFT
    )

    # ── Subtitle ──────────────────────────────────────────────────
    tb_sub = None
    if subtitle_text:
        tb_sub = _add_textbox(
            slide, subtitle_text,
            0.55, 4.15, 10.0, 1.0,
            Pt(20), False, tmpl['accent'], tmpl['font_name'], PP_ALIGN.LEFT
        )

    # ── Bottom tag (brand + date) ─────────────────────────────────
    _add_textbox(
        slide, 'Powered by SlideAI',
        0.55, 6.9, 4.0, 0.4,
        Pt(11), False, '#FFFFFF', tmpl['font_name'], PP_ALIGN.LEFT, italic=True
    )

    # ── Animations ────────────────────────────────────────────────
    anim = [(tb_title.shape_id, 300)]
    if tb_sub:
        anim.append((tb_sub.shape_id, 800))
    anim.append((line.shape_id, 600))
    _build_animations(slide, anim)
    _add_transition(slide._element, 'fade')
    return slide


def _make_content_slide(prs, slide_data, tmpl, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_text = slide_data.get('title', '')
    bullets    = slide_data.get('bullet_points', [])

    # ── Background ────────────────────────────────────────────────
    _add_rect(slide, 0, 0, 13.33, 7.5, tmpl['background'])

    # ── Header bar ────────────────────────────────────────────────
    header = _add_rect(slide, 0, 0, 13.33, 1.1, tmpl['primary'])
    accent_stripe = _add_rect(slide, 0, 1.1, 13.33, 0.06, tmpl['accent'])

    # ── Title in header ───────────────────────────────────────────
    tb_title = _add_textbox(
        slide, title_text,
        0.4, 0.16, 12.2, 0.88,
        Pt(26), True, '#FFFFFF', tmpl['font_name']
    )

    # ── Numbered bullet points with badge circles ─────────────────
    anim = [(header.shape_id, 0), (tb_title.shape_id, 150)]
    y = 1.35
    for i, bullet in enumerate(bullets[:6]):
        delay = 350 + i * 200

        # Number badge circle
        badge_color = tmpl['accent'] if i % 2 == 0 else tmpl['secondary']
        badge, badge_tb = _add_number_badge(
            slide, i + 1, 0.35, y + 0.04, badge_color, '#FFFFFF'
        )

        # Bullet text box
        bullet_box = slide.shapes.add_textbox(
            Inches(0.9), Inches(y), Inches(12.0), Inches(0.72)
        )
        tf = bullet_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = bullet
        _set_font(run, Pt(17), False, tmpl['text_color'], tmpl['font_name'])

        anim.append((badge.shape_id, delay))
        anim.append((bullet_box.shape_id, delay + 60))
        y += 0.9

    # ── Slide number ──────────────────────────────────────────────
    _add_slide_number(slide, slide_num, total, tmpl)

    _build_animations(slide, anim)
    _add_transition(slide._element, 'push')
    return slide


def _make_chart_slide(prs, slide_data, tmpl, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_text = slide_data.get('title', '')
    bullets    = (slide_data.get('bullet_points') or [])[:4]

    raw_cd         = slide_data.get('chart_data') or {}
    raw_cd['chart_type'] = slide_data.get('chart_type', 'bar')

    # ── Background ────────────────────────────────────────────────
    _add_rect(slide, 0, 0, 13.33, 7.5, tmpl['background'])

    # ── Header bar ────────────────────────────────────────────────
    header = _add_rect(slide, 0, 0, 13.33, 1.05, tmpl['primary'])
    _add_rect(slide, 0, 1.05, 13.33, 0.06, tmpl['accent'])

    # ── Title ─────────────────────────────────────────────────────
    tb_title = _add_textbox(
        slide, title_text,
        0.4, 0.14, 12.2, 0.84,
        Pt(26), True, '#FFFFFF', tmpl['font_name']
    )

    # ── Chart + optional insight panel ────────────────────────────
    anim = [(header.shape_id, 0), (tb_title.shape_id, 150)]

    if bullets:
        chart_gf = _add_chart(slide, raw_cd, 0.35, 1.22, 8.8, 5.9, tmpl)
        anim.append((chart_gf.shape_id, 350))

        # Right insight panel background
        panel = _add_rect(slide, 9.3, 1.22, 3.8, 5.9, tmpl['primary'], rounded=True)
        anim.append((panel.shape_id, 700))

        # Panel header
        ph = _add_textbox(slide, 'KEY INSIGHTS', 9.5, 1.35, 3.4, 0.4,
                          Pt(11), True, tmpl['accent'], tmpl['font_name'])
        anim.append((ph.shape_id, 750))

        for i, b in enumerate(bullets):
            delay = 850 + i * 200
            dot = _add_oval(slide, 9.5, 2.0 + i * 1.08, 0.18, 0.18, tmpl['accent'])
            insight_tb = slide.shapes.add_textbox(
                Inches(9.82), Inches(1.95 + i * 1.08), Inches(3.1), Inches(0.88)
            )
            tf = insight_tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = b
            _set_font(run, Pt(13), False, '#FFFFFF', tmpl['font_name'])
            anim += [(dot.shape_id, delay), (insight_tb.shape_id, delay + 50)]
    else:
        chart_gf = _add_chart(slide, raw_cd, 0.5, 1.22, 12.2, 5.9, tmpl)
        anim.append((chart_gf.shape_id, 350))

    _add_slide_number(slide, slide_num, total, tmpl)
    _build_animations(slide, anim)
    _add_transition(slide._element, 'wipe')
    return slide


def _make_two_col_slide(prs, slide_data, tmpl, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_text = slide_data.get('title', '')
    bullets    = slide_data.get('bullet_points', [])
    mid        = max(1, len(bullets) // 2)
    left_bul   = bullets[:mid]
    right_bul  = bullets[mid:]

    # ── Background ────────────────────────────────────────────────
    _add_rect(slide, 0, 0, 13.33, 7.5, tmpl['background'])

    # ── Header ────────────────────────────────────────────────────
    header = _add_rect(slide, 0, 0, 13.33, 1.05, tmpl['primary'])
    _add_rect(slide, 0, 1.05, 13.33, 0.06, tmpl['accent'])

    tb_title = _add_textbox(
        slide, title_text,
        0.4, 0.14, 12.2, 0.84,
        Pt(26), True, '#FFFFFF', tmpl['font_name']
    )

    # ── Column header backgrounds ─────────────────────────────────
    left_hdr  = _add_rect(slide, 0.35, 1.25, 5.9, 0.5, tmpl['secondary'], rounded=True)
    right_hdr = _add_rect(slide, 7.1, 1.25, 5.9, 0.5, tmpl['secondary'], rounded=True)

    _add_textbox(slide, 'COLUMN A', 0.35, 1.25, 5.9, 0.5,
                 Pt(12), True, '#FFFFFF', tmpl['font_name'], PP_ALIGN.CENTER)
    _add_textbox(slide, 'COLUMN B', 7.1, 1.25, 5.9, 0.5,
                 Pt(12), True, '#FFFFFF', tmpl['font_name'], PP_ALIGN.CENTER)

    # ── Vertical divider ──────────────────────────────────────────
    divider = _add_rect(slide, 6.58, 1.2, 0.07, 5.9, tmpl['accent'])

    # ── Left bullets ──────────────────────────────────────────────
    anim = [(header.shape_id, 0), (tb_title.shape_id, 150),
            (left_hdr.shape_id, 300), (right_hdr.shape_id, 350), (divider.shape_id, 400)]

    for i, b in enumerate(left_bul[:4]):
        delay = 500 + i * 200
        badge_c = _add_oval(slide, 0.38, 2.05 + i * 1.0, 0.34, 0.34, tmpl['accent'])
        tb = slide.shapes.add_textbox(
            Inches(0.85), Inches(2.02 + i * 1.0), Inches(5.55), Inches(0.85)
        )
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; run = p.add_run(); run.text = b
        _set_font(run, Pt(15), False, tmpl['text_color'], tmpl['font_name'])
        anim += [(badge_c.shape_id, delay), (tb.shape_id, delay + 50)]

    for i, b in enumerate(right_bul[:4]):
        delay = 550 + i * 200
        badge_c = _add_oval(slide, 7.12, 2.05 + i * 1.0, 0.34, 0.34, tmpl['primary'])
        tb = slide.shapes.add_textbox(
            Inches(7.6), Inches(2.02 + i * 1.0), Inches(5.3), Inches(0.85)
        )
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; run = p.add_run(); run.text = b
        _set_font(run, Pt(15), False, tmpl['text_color'], tmpl['font_name'])
        anim += [(badge_c.shape_id, delay), (tb.shape_id, delay + 50)]

    _add_slide_number(slide, slide_num, total, tmpl)
    _build_animations(slide, anim)
    _add_transition(slide._element, 'push')
    return slide


def _make_divider_slide(prs, slide_data, tmpl, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_text    = slide_data.get('title', '')
    subtitle_text = slide_data.get('subtitle', '')

    # ── Full primary background ───────────────────────────────────
    _add_rect(slide, 0, 0, 13.33, 7.5, tmpl['primary'])

    # ── Decorative circles ────────────────────────────────────────
    _add_oval(slide, -1.5, -1.5, 5.0, 5.0, tmpl['secondary'])
    _add_oval(slide, 10.3, 4.5,  4.0, 4.0, tmpl['secondary'])

    # ── Section number (large, background) ───────────────────────
    sec_num_str = str(slide_num)
    _add_textbox(
        slide, sec_num_str,
        0.5, 0.5, 3.0, 3.5,
        Pt(140), True, tmpl['secondary'], tmpl['font_name'], PP_ALIGN.LEFT
    )

    # ── Horizontal accent line ────────────────────────────────────
    line = _add_rect(slide, 0.55, 4.1, 5.5, 0.08, tmpl['accent'])

    # ── Title ─────────────────────────────────────────────────────
    tb_title = _add_textbox(
        slide, title_text,
        0.55, 2.3, 11.5, 2.0,
        Pt(40), True, '#FFFFFF', tmpl['font_name'], PP_ALIGN.LEFT
    )

    # ── Subtitle ──────────────────────────────────────────────────
    tb_sub = None
    if subtitle_text:
        tb_sub = _add_textbox(
            slide, subtitle_text,
            0.55, 4.4, 10.5, 0.9,
            Pt(18), False, tmpl['accent'], tmpl['font_name'], PP_ALIGN.LEFT
        )

    anim = [(line.shape_id, 200), (tb_title.shape_id, 400)]
    if tb_sub:
        anim.append((tb_sub.shape_id, 900))
    _build_animations(slide, anim)
    _add_transition(slide._element, 'zoom')
    return slide


def _make_conclusion_slide(prs, slide_data, tmpl, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_text = slide_data.get('title', 'Key Takeaways')
    bullets    = slide_data.get('bullet_points', [])

    # ── Background ────────────────────────────────────────────────
    _add_rect(slide, 0, 0, 13.33, 7.5, tmpl['background'])

    # ── Header bar ────────────────────────────────────────────────
    header = _add_rect(slide, 0, 0, 13.33, 1.05, tmpl['primary'])
    _add_rect(slide, 0, 1.05, 13.33, 0.06, tmpl['accent'])

    # ── "CONCLUSION" label ────────────────────────────────────────
    _add_textbox(slide, 'CONCLUSION', 0.4, 0.04, 3.5, 0.38,
                 Pt(11), True, tmpl['accent'], tmpl['font_name'])

    # ── Title ─────────────────────────────────────────────────────
    tb_title = _add_textbox(
        slide, title_text,
        0.4, 0.38, 12.2, 0.6,
        Pt(24), True, '#FFFFFF', tmpl['font_name']
    )

    anim = [(header.shape_id, 0), (tb_title.shape_id, 150)]

    # ── Takeaway cards (up to 5) ──────────────────────────────────
    is_two_col = len(bullets) > 3
    cards_per_col = 3
    col_configs = []

    if is_two_col:
        col1 = bullets[:3]
        col2 = bullets[3:6]
        col_configs = [
            (col1, 0.35,  1.3),
            (col2, 6.85, 1.3),
        ]
    else:
        col_configs = [(bullets[:5], 0.35, 1.3)]

    for col_list, col_x, col_y in col_configs:
        for i, b in enumerate(col_list):
            delay = 350 + i * 220
            card_w = 6.0 if is_two_col else 12.5

            # Card background
            card = _add_rect(slide, col_x, col_y + i * 1.85, card_w, 1.62,
                              tmpl['primary'] if i % 2 == 0 else tmpl['secondary'],
                              rounded=True)

            # Left accent bar
            acc = _add_rect(slide, col_x, col_y + i * 1.85, 0.12, 1.62, tmpl['accent'])

            # Check mark
            chk = _add_textbox(
                slide, '✓', col_x + 0.22, col_y + i * 1.85 + 0.5, 0.5, 0.6,
                Pt(22), True, tmpl['accent'], tmpl['font_name'], PP_ALIGN.CENTER
            )

            # Card text
            tb = slide.shapes.add_textbox(
                Inches(col_x + 0.75), Inches(col_y + i * 1.85 + 0.2),
                Inches(card_w - 0.85), Inches(1.25)
            )
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; run = p.add_run(); run.text = b
            _set_font(run, Pt(15), False, '#FFFFFF', tmpl['font_name'])

            anim += [(card.shape_id, delay), (tb.shape_id, delay + 80)]

    _add_slide_number(slide, slide_num, total, tmpl)
    _build_animations(slide, anim)
    _add_transition(slide._element, 'fade')
    return slide


# ═══════════════════════ DISPATCH ════════════════════════════════

_BUILDERS = {
    'title':      _make_title_slide,
    'content':    _make_content_slide,
    'chart':      _make_chart_slide,
    'two_col':    _make_two_col_slide,
    'divider':    _make_divider_slide,
    'conclusion': _make_conclusion_slide,
}


def build_pptx(slides_data, template_id, output_dir):
    tmpl  = get_template(template_id)
    prs   = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    total = len(slides_data)
    for idx, slide_data in enumerate(slides_data, start=1):
        stype   = slide_data.get('slide_type', 'content')
        builder = _BUILDERS.get(stype, _make_content_slide)
        builder(prs, slide_data, tmpl, idx, total)

    filename = f"{uuid.uuid4().hex[:12]}_presentation.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    return filepath, filename
