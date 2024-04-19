import time
import uuid
from flask import Flask, render_template, request, send_file
from pptx import Presentation
import openai
import os


from dotenv import load_dotenv
load_dotenv()


app = Flask(__name__)

# Replace 'your-api-key' with your actual OpenAI API key
openai.api_key = os.getenv("OPENAPI_KEY")

@app.route('/', methods=['GET', 'POST'])
def home():
    if request.method == 'POST':
        num_slides = int(request.form.get('num_slides', 10))  # Or any other way you get this number
        prompt = request.form.get('prompt')
        print("Prompt is:", prompt)
        
        # Use OpenAI to generate outlines for the specified number of slides
        outlines = generate_outlines(prompt, num_slides)
        
        # For each outline, generate content
        slides_content = [generate_content_for_outline(outline) for outline in outlines]
        
        # Create a presentation with the generated content
        ppt = create_presentation(slides_content)
        
        # Save the presentation
        ppt_file = f"generated_presentation_{uuid.uuid4().hex[:8]}.pptx"
        ppt.save(ppt_file)

        return send_file(ppt_file, as_attachment=True)

    return render_template('index.html')
def generate_outlines(main_topic, num_slides):
    # Craft a prompt to ask for a list of related subtopics or sections for the main topic
    prompt = f"Give me {num_slides} slides possible outlines for topic {main_topic} for Powerpoint presentation, including an introduction and conclusion. Length of the outlines must be shorter as one line heading."

    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        temperature=0.5,
        max_tokens=200,  # Adjust as needed
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )

    outlines = response.choices[0].text.strip().split('\n')
    # Ensure that the number of outlines matches num_slides, trim or fill if necessary
    return outlines[:num_slides] if len(outlines) >= num_slides else outlines + ["Additional Content" for _ in range(num_slides - len(outlines))]

def generate_content_for_outline(outline):
    # Craft a prompt to generate detailed content for each outline
    prompt = f"Explain the topic '{outline}' in detail suitable for a single slide presentation."

    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        temperature=0.5,
        max_tokens=150,  # Adjust as needed
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )

    content = response.choices[0].text.strip()
    return {"title": outline, "content": content}


def create_presentation(slides_content):
    ppt = Presentation()
    
    # Clear all existing slides, if necessary
    #while len(ppt.slides) > 0:
    #    xml_slides = ppt.slides._sldIdLst
    #    slides = list(xml_slides)
    #    xml_slides.remove(slides[0])

    for slide_content in slides_content:
        slide_layout = ppt.slide_layouts[1]  # Title and Content layout
        slide = ppt.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_content["title"]
        slide.placeholders[1].text = slide_content["content"]
    return ppt


if __name__ == '__main__':
    app.run(debug=True)
