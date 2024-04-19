import os
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
import openai

from dotenv import load_dotenv
load_dotenv()

# Initialize your OpenAI API key
openai.api_key = os.getenv("OPENAPI_KEY")

def generate_slide_titles(topic, num_slides):
    prompt = f"Create a list of titles for a {num_slides}-slide PowerPoint presentation on the topic: '{topic}'."
    response = openai.Completion.create(
        engine="gpt-3.5-turbo-instruct",
        prompt=prompt,
        max_tokens=50 * num_slides
    )
    titles = response.choices[0].text.strip().split('\n')
    return titles

def generate_slide_content(title):
    prompt = f"Generate a concise list of bullet points for a PowerPoint slide with the title '{title}'."
    response = openai.Completion.create(
        engine="gpt-3.5-turbo-instruct",
        prompt=prompt,
        max_tokens=200
    )
    content = response.choices[0].text.strip()
    return content


def add_slide_to_presentation(prs, title, content):
    slide_layout = prs.slide_layouts[5]  # Choosing a slide layout that is supposed to have a title and content.
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    # Instead of using placeholders, add a text box directly
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    # Split content into lines and add to text frame
    for point in content.split('\n'):
        p = tf.add_paragraph()
        p.text = point
        p.level = 0  # Adjust level as necessary for sub-points
        p.font.size = Pt(12)  # Adjust font size as necessary

def count_pptx_files(directory):
    # List all files in the directory
    files = os.listdir(directory)
    # Filter out all files with the .pptx extension
    pptx_files = [file for file in files if file.endswith('.pptx')]
    # Return the count of .pptx files
    return len(pptx_files)


def main():
    topic = input("Enter the topic for the PowerPoint presentation: ")
    num_slides = int(input("How many slides do you want to generate? "))

    prs = Presentation()

    print(f"Generating an outline for {num_slides} slides on the topic '{topic}'...")
    slide_titles = generate_slide_titles(topic, num_slides)
    
    for i, title in enumerate(slide_titles, start=1):
        print(f"Slide {i}: {title}")

    generate_content = input("Do you want to generate detailed content for each slide? (Y/N) ").strip().lower()
    
    if generate_content == 'y':
        for i, title in enumerate(slide_titles, start=1):
            print(f"\nGenerating content for Slide {i}: {title}...")
            content = generate_slide_content(title)
            add_slide_to_presentation(prs, title, content)
            print(f"Content for Slide {i}:\n{content}")

    pptx_file_name = f"output/{count_pptx_files('output')+1}_presentation.pptx"
    prs.save(pptx_file_name)
    print(f"Presentation saved as {pptx_file_name}")

if __name__ == "__main__":
    main()
