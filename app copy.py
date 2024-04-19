from flask import Flask, render_template, request, send_file
from pptx import Presentation
import openai
import os

from dotenv import load_dotenv
load_dotenv()

app = Flask(__name__)

# Replace 'your-api-key' with your actual OpenAI API key
openai.api_key = os.getenv("OPENAPI_KEY")

@app.route('/', methods=['GET', 'POST'])
def home():
    if request.method == 'POST':
        num_slides = int(request.form.get('num_slides', 10))  # Default to 5 slides

        prompt = request.form.get('prompt', "Powerpoint Presentation with AI")
        
        # Use OpenAI to generate sections based on the prompt
        sections = generate_sections_with_openai(prompt)
        
        # Create a presentation
        ppt = create_presentation(sections)
        
        # Save the presentation
        ppt_file = "generated_presentation.pptx"
        ppt.save(ppt_file)

        return send_file(ppt_file, as_attachment=True)

    return render_template('index.html')

def generate_sections_with_openai(prompt):
    structured_prompt = (
        f"Create a detailed outline for a presentation on {prompt}. "
        "Include sections for an introduction, benefits, and conclusion. "
        "Provide a brief description for each section and suggest images that could accompany the text. "
        "Format the content with titles for each section and bullet points for details."
    )
    response = openai.Completion.create(
        engine="text-davinci-003",  # or another appropriate engine
        prompt=structured_prompt,  # Your structured prompt to the model
        temperature=0.5,
        max_tokens=1000,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )
    
    text = response.choices[0].text.strip()
    sections = parse_sections(text)
    return sections

def parse_sections(text):
    # Implement your logic to parse the text into title/content for each section.
    # This is a simplistic approach; you might need a more sophisticated parser.
    sections = []
    for part in text.split("\n\n"):
        if part.strip():
            title, *content = part.strip().split("\n")
            sections.append({"title": title, "content": "\n".join(content)})
    return sections

def create_presentation(sections):
    ppt = Presentation()
    for section in sections:
        slide_layout = ppt.slide_layouts[1]  # Title and Content
        slide = ppt.slides.add_slide(slide_layout)
        slide.shapes.title.text = section["title"]
        slide.placeholders[1].text = section["content"]
    return ppt

if __name__ == '__main__':
    app.run(debug=True)
