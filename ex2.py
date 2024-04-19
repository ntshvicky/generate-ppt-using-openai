from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches
import requests

from dotenv import load_dotenv
load_dotenv()

def html_to_ppt(html_url):
    # Fetch the HTML content
    response = requests.get(html_url)
    soup = BeautifulSoup(response.text, 'html.parser')

    # Create a PowerPoint presentation
    prs = Presentation()

    # Assuming a simple case where we just take all 'p' and 'img' tags
    for element in soup.find_all(['p', 'img']):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Using a blank slide layout

        if element.name == 'p':
            # Add a textbox and put the text there
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8.5), Inches(1))
            tf = txBox.text_frame
            tf.text = element.get_text()

        elif element.name == 'img':
            # Download the image and add it
            img_url = "https://nitishsrivastava.in/{}".format(element['src']) if 'https' not in element['src'] else element['src']
            img_response = requests.get(img_url)
            img_path = 'temp_image.jpg'
            with open(img_path, 'wb') as f:
                f.write(img_response.content)
            slide.shapes.add_picture(img_path, Inches(1), Inches(2))

    prs.save('output.pptx')

html_to_ppt('https://nitishsrivastava.in')
