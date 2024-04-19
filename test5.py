import os
import pptx
from pptx.util import Inches
import requests
from io import BytesIO
from tqdm import tqdm

from dotenv import load_dotenv
load_dotenv()


from openai import OpenAI

# Initialize your OpenAI API key

client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY", os.getenv("OPENAPI_KEY")))


def generate_complete_slide_content(topic, num_slides):

    content_prompt = f'''Generate content for a PowerPoint presentation on "{topic}". 
                    
                    The presentation should consist of {num_slides} slides with the following details:
                    
                    Preferred Structure:
                        Slide 1: Introduction(add small introduction on topic)
                        Slide 2-{num_slides-2}: Main Content (Different aspects of the impact)
                        Slide {num_slides-1}: Case Studies or Notable Examples
                        Slide {num_slides}: Conclusion
                    
                    Slide Details: For each slide in the PowerPoint presentation provide a concise bullet point followed by a brief explanation of 2 to 3 sentences. 
                    Ensure that each bullet point is clear and informative, covering key aspects of {topic}. 
                    The explanations should elaborate on the bullet point, providing context or additional details suitable for a professional audience. 
                    Each slide should contain 5 to 6 of these bullet points with explanations, making the content informative and engaging.
                    
                    The total content should contain more than 1000 but less than 1200 tokens.
                        '''


    # example with a system message
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "user", "content": content_prompt},
        ],
        temperature=0,
    )
    #print(content_prompt)
    return response.choices[0].message.content


def create_presentation(content, topic):
    prs = pptx.Presentation()

    # Splitting the content into slides
    slides_content = content.split("\n\n")

    # Initialize progress bar
    progress_bar = tqdm(total=len(slides_content), desc="Generating Slides")

    for slide_content in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title, body = slide.shapes.title, slide.placeholders[1]
        title.text = f"{topic}"
        body.text = slide_content

        # Generate and add an image to the slide
        '''
        image_url = generate_image(f"Image for powerpoint slides, related to topic '{slide_content}'. It should be simple related, and memorable.")
        if image_url:
            image_stream = download_image(image_url)
            if image_stream:
                slide.shapes.add_picture(image_stream, Inches(2), Inches(2), width=Inches(4), height=Inches(3))
        '''
        # Update progress bar after each slide
        progress_bar.update(1)

    # Close the progress bar
    progress_bar.close()

    return prs

def main():
    topic = input("Enter a topic for the PowerPoint presentation: ")
    num_slides = int(input("Enter the number of slides you want: "))
    content = generate_complete_slide_content(topic, num_slides)
    if content:
        presentation = create_presentation(content, topic)
        file_name = f"{topic}_presentation.pptx"
        presentation.save(file_name)
        print(f"Presentation saved as {file_name}")

if __name__ == "__main__":
    main()
