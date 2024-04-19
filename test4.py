import os
import re
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
import openai

from dotenv import load_dotenv
load_dotenv()

# Initialize your OpenAI API key
openai.api_key = os.getenv("OPENAPI_KEY")

def clean_title(title):
    # Remove leading numbers, periods, and spaces
    title = re.sub(r'^\d+\.\s*', '', title)

    # Remove double quotes
    title = title.replace('"', '')

    return title

def select_category():
    categories = {
        1: "Educational",
        2: "Business",
        3: "Technical",
        4: "Sales",
        5: "Marketing",
        6: "Motivational",
        7: "Training",
        8: "Financial",
        9: "Proposal",
        10: "Project Update or Status",
        11: "Storytelling",
        12: "Webinars",
        13: "Interactive",
        14: "Informative"
    }
    print("Select a category for your presentation:")
    for key, value in categories.items():
        print(f"{key}: {value}")
    choice = int(input("Enter the number corresponding to your choice: "))
    return categories.get(choice, "General")

def generate_slide_titles(topic, num_slides, category):
    prompt = f"Create a list of titles for a {num_slides}-slide {category} PowerPoint presentation on the topic: '{topic}'."
    response = openai.Completion.create(
        engine="gpt-3.5-turbo-instruct",
        prompt=prompt,
        max_tokens=50 * num_slides,
        temperature=0,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )
    titles = response.choices[0].text.strip().split('\n')
    return titles

'''
def generate_complete_slide_content(title):
    # Initial prompt for the first part of the content
    prompt = f"Generate a detailed, concise list of bullet points for the first part of a PowerPoint slide with the title '{title}'. Each bullet point should provide specific information, suitable for graphical representation."

    # Initialize variables
    content = ""
    part = 1
    all_content_received = False

    while not all_content_received:
        response = openai.Completion.create(
            engine="gpt-3.5-turbo-instruct",
            prompt=prompt,
            max_tokens=200
        )
        part_content = response.choices[0].text.strip()

        # Check if the received content is less than a certain threshold, indicating end of content
        if len(part_content.split()) < 50:  # Threshold can be adjusted
            all_content_received = True

        content += part_content

        # Update prompt for the next part
        part += 1
        prompt = f"Generate the next part of the list of bullet points for a PowerPoint slide with the title '{title}', continuing from the last point."

    return content
'''

def generate_slide_content(title, category):
    prompt = f"Generate a brief and concise list of bullet points for a {category.lower()} presentation that explain the '{title}'. Each bullet point should provide key information, be direct and to-the-point, suitable for a slide format, include applications, benefits, challenges, and future prospects related to the topic. Content should be engaging and tailored to the audience."
    response = openai.Completion.create(
        engine="gpt-3.5-turbo-instruct",
        prompt=prompt,
        max_tokens=256,
        temperature=0,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )
    content = response.choices[0].text.strip()
    return content

def add_slide_to_presentation(prs, title, content):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = clean_title(title)

    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = True

    # Split content into paragraphs and subpoints
    for line in content.split('\n'):
        if line.startswith('- '):  # Subpoint indicator
            p = tf.add_paragraph()
            p.text = line[2:]  # Remove '- ' from beginning
            p.level = 1  # Subpoint level
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0  # Main point level
        p.font.size = Pt(12)

def count_pptx_files(directory):
    files = os.listdir(directory)
    pptx_files = [file for file in files if file.endswith('.pptx')]
    return len(pptx_files)

def main():
    topic = input("Enter the topic for the PowerPoint presentation: ")
    category = select_category()

    num_slides = int(input("How many slides do you want to generate? "))

    prs = Presentation()

    print(f"Generating an outline for {num_slides} slides on the topic '{topic}' in the {category} category...")
    slide_titles = generate_slide_titles(topic, num_slides, category)
    
    for i, title in enumerate(slide_titles, start=1):
        print(f"Slide {i}: {title}")

    generate_content = input("Do you want to generate detailed content for each slide? (Y/N) ").strip().lower()
    
    if generate_content == 'y':
        for i, title in enumerate(slide_titles, start=1):
            print(f"\nGenerating content for Slide {i}: {title}...")
            content = generate_slide_content(title, category)
            add_slide_to_presentation(prs, title, content)
            print(f"Content for Slide {i}:\n{content}")

        pptx_file_name = f"output/{count_pptx_files('output')+1}_presentation.pptx"
        prs.save(pptx_file_name)
        print(f"Presentation saved as {pptx_file_name}")
        return
    else:
        print('Generating slides cancelled.')

if __name__ == "__main__":
    main()