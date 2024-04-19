from flask import Flask, render_template, request, send_file
from pptx import Presentation
import openai
import os

from dotenv import load_dotenv
load_dotenv()

app = Flask(__name__)

# Replace 'your-api-key' with your actual OpenAI API key
openai.api_key = os.getenv("OPENAPI_KEY")

@app.route('/', methods=['GET', 'POST'])
def home():
    if request.method == 'POST':
        prompt = request.form.get('prompt', 'Default Prompt')
        
        # Use OpenAI to generate sections based on the prompt
        sections = generate_sections_with_openai(prompt)
        
        # Create a presentation
        ppt = create_presentation(sections)
        
        # Save the presentation
        ppt_file = "generated_presentation.pptx"
        ppt.save(ppt_file)

        return send_file(ppt_file, as_attachment=True)

    return render_template('index.html')

def generate_sections_with_openai(prompt):
    response = openai.Completion.create(
      engine="text-davinci-003", # or another appropriate engine
      prompt=prompt, # Your prompt to the model
      temperature=0.7,
      max_tokens=500,
      top_p=1,
      frequency_penalty=0,
      presence_penalty=0
    )
    
    text = response.choices[0].text.strip()
    # Now, you'd parse the text into sections. This might involve more complex logic
    # depending on how you structure your prompts and what you expect back.
    sections = parse_sections(text)
    return sections

def parse_sections(text):
    # Implement your logic to parse the text into title/content for each section.
    # For now, let's assume it's a simple split for illustration.
    parts = text.split("\n\n")
    sections = []
    for part in parts:
        title, content = part.split("\n")[0], "\n".join(part.split("\n")[1:])
        sections.append({"title": title, "content": content})
    return sections

def create_presentation(sections):
    ppt = Presentation()
    for section in sections:
        slide_layout = ppt.slide_layouts[1]  # Title and Content
        slide = ppt.slides.add_slide(slide_layout)
        slide.shapes.title.text = section["title"]
        slide.placeholders[1].text = section["content"]
    return ppt

if __name__ == '__main__':
    app.run(debug=True)
